#!/usr/bin/env python3
"""
Genererer en PowerPoint-presentasjon om Vibe Koding
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Aftenbladet blå farge
AFTENBLAD_BLUE = RGBColor(0, 51, 102)  # Mørk blå
AFTENBLAD_LIGHT_BLUE = RGBColor(0, 102, 153)
WHITE = RGBColor(255, 255, 255)

def add_expandable_slide(prs, title, intro_lines, expandable_title, expandable_lines):
    """Legger til en slide med en "utvidbar" seksjon (visuelt markert)."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, AFTENBLAD_BLUE)

    # Tittel
    left = Inches(0.5)
    top = Inches(0.3)
    width = Inches(9)
    height = Inches(1)

    title_box = slide.shapes.add_textbox(left, top, width, height)
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Intro
    content_top = Inches(1.5)
    content_height = Inches(2.2)
    content_box = slide.shapes.add_textbox(left, content_top, width, content_height)
    tf = content_box.text_frame
    tf.word_wrap = True

    for i, line in enumerate(intro_lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {line}"
        p.font.size = Pt(20)
        p.font.color.rgb = WHITE
        p.space_after = Pt(10)

    # "Expandable" panel
    panel_top = Inches(3.9)
    panel_height = Inches(2.9)
    panel = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left,
        panel_top,
        width,
        panel_height,
    )
    panel.fill.solid()
    panel.fill.fore_color.rgb = AFTENBLAD_LIGHT_BLUE
    panel.line.color.rgb = RGBColor(173, 216, 230)
    panel.line.width = Pt(1.25)

    tf = panel.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.25)
    tf.margin_right = Inches(0.25)
    tf.margin_top = Inches(0.15)
    tf.vertical_anchor = MSO_ANCHOR.TOP

    p = tf.paragraphs[0]
    p.text = f"▶ {expandable_title}"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.space_after = Pt(10)

    for line in expandable_lines:
        p = tf.add_paragraph()
        p.text = f"• {line}"
        p.font.size = Pt(16)
        p.font.color.rgb = WHITE
        p.space_after = Pt(6)

    return slide

def set_slide_background(slide, color):
    """Setter bakgrunnsfarge på en slide"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_title_slide(prs, title, subtitle=""):
    """Legger til en tittelslide"""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, AFTENBLAD_BLUE)
    
    # Tittel
    left = Inches(0.5)
    top = Inches(2.5)
    width = Inches(9)
    height = Inches(1.5)
    
    title_box = slide.shapes.add_textbox(left, top, width, height)
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    if subtitle:
        # Undertittel
        sub_top = Inches(4.2)
        sub_box = slide.shapes.add_textbox(left, sub_top, width, Inches(1))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(24)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, bullet_points, has_link=None):
    """Legger til en innholdsslide med punktliste"""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, AFTENBLAD_BLUE)
    
    # Tittel
    left = Inches(0.5)
    top = Inches(0.3)
    width = Inches(9)
    height = Inches(1)
    
    title_box = slide.shapes.add_textbox(left, top, width, height)
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Innhold
    content_top = Inches(1.5)
    content_height = Inches(5)
    content_box = slide.shapes.add_textbox(left, content_top, width, content_height)
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, point in enumerate(bullet_points):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {point}"
        p.font.size = Pt(20)
        p.font.color.rgb = WHITE
        p.space_after = Pt(12)
    
    if has_link:
        p = tf.add_paragraph()
        p.text = f"\n🔗 {has_link}"
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(173, 216, 230)  # Lysblå for lenker
    
    return slide

def add_section_slide(prs, title):
    """Legger til en seksjonsslide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, AFTENBLAD_LIGHT_BLUE)
    
    left = Inches(0.5)
    top = Inches(3)
    width = Inches(9)
    height = Inches(1.5)
    
    title_box = slide.shapes.add_textbox(left, top, width, height)
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_meeting_slide(prs, title, meetings):
    """Legger til en slide med møteoversikt"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, AFTENBLAD_BLUE)
    
    # Tittel
    left = Inches(0.3)
    top = Inches(0.2)
    width = Inches(9.4)
    height = Inches(0.8)
    
    title_box = slide.shapes.add_textbox(left, top, width, height)
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Møter
    content_top = Inches(1)
    content_box = slide.shapes.add_textbox(left, content_top, width, Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, meeting in enumerate(meetings):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = meeting
        p.font.size = Pt(14)
        p.font.color.rgb = WHITE
        p.space_after = Pt(4)
    
    return slide

def create_presentation():
    """Hovedfunksjon som lager presentasjonen"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # 1. Tittelslide
    add_title_slide(
        prs,
        "Vibe koding",
        "Intro · eksempler · automasjon · bots\nDatadrevet journalistikk i Stavanger Aftenblad\n26. januar 2026",
    )
    
    # 2. Intro: hva mener vi med vibe koding?
    add_content_slide(prs, "Hva mener vi med «vibe koding» for journalister?", [
        "Små verktøy som sparer tid i hverdagen",
        "Data, grafikk og automatisering som styrker journalistikken",
        "Neste side: konkrete eksempler med forklaring"
    ])

    # 3. Hva er vibe koding?
    add_expandable_slide(
        prs,
        "Hva er «vibe coding»?",
        [
            "Konseptet med vibbekoding er enkelt:",
            "Du lar AI-en kjøre på, og kode alt for deg."
        ],
        "Ulemper",
        [
            "Du kan få kode du ikke helt forstår",
            "Bugs, sikkerhetshull og teknisk gjeld skjules lettere",
            "Mer tid går til feilsøking hvis du ikke tester",
            "Ujevn kvalitet uten god struktur og review"
        ],
    )
    
    # 4. Seksjonsslide - Prosjekter
    add_section_slide(prs, "📊 Eksempler på prosjekter")
    
    # 4. Forsinkede busser
    add_content_slide(prs, "Forsinkede busser", [
        "Utviklet av: Andreas Askildsen",
        "Datakilde: Entur API",
        "Sanntidsdata på bussavganger",
        "Visualisering av forsinkelser",
        "Interaktiv løsning for leserne"
    ], "aftenbladet-editorial-sa-entur-tij.pages.schibsted.ghe.com")
    
    # 5. Restauranter
    add_content_slide(prs, "Restaurantoversikter", [
        "Utviklet av: Eilin Lindvoll (Sandnes/Stavanger)",
        "Datakilde: Brønnøysundregistrene",
        "Automatisk oppdatert database",
        "Dekker: Sandnes, Stavanger, Jæren, Gjesdal",
        "Inline HTML – publiseres direkte i saker"
    ])
    
    # 6. Flere interaktive prosjekter
    add_content_slide(prs, "Flere interaktive løsninger", [
        "Oversitting på Vassøy-ferja (Kolumbus-data)",
        "Ryfast-takster 2026",
        "Beste akebakke (Karl Almelid m.fl.)",
        "Bygninger i strandsonen (telleverk)",
        "Stavanger kommunes eiendommer (Elisabeth Risa)"
    ])
    
    # 7. Seksjonsslide - Automatisering
    add_section_slide(prs, "🤖 Automatisering og bots")
    
    # 8. Slack-varsling
    add_content_slide(prs, "Slack-varsling", [
        "Varsler når konkurrenter publiserer nyheter",
        "Automatisk overvåkning",
        "Tidsbesparelser for redaksjonen",
        "Holder oss oppdatert på nyhetsbildet"
    ])
    
    # 9. Dagsorden-bot
    add_content_slide(prs, "Dagsorden-bot", [
        "Skraping av møtekalendere",
        "Kommuner, fylkeskommune, m.fl.",
        "Automatisk Slack-varsling",
        "39 møter de neste 10 dagene (Ryfylke, Dalane)",
        "Aldri gå glipp av et viktig møte"
    ])
    
    # 10. Møteoversikt eksempel
    add_meeting_slide(prs, "Eksempel: Politiske møter denne uken", [
        "📅 Mandag 26. januar:",
        "   • Eldreråd (Lund kommune) - kl. 08:30",
        "   • Kommunestyret (Eigersund) - kl. 18:00",
        "",
        "📅 Tirsdag 27. januar:",
        "   • Opplæringsutvalget (Rogaland fylkeskommune) - kl. 11:30",
        "   • Formannskapet (Eigersund) - kl. 13:01",
        "",
        "📅 Onsdag 28. januar:",
        "   • Samferdselsutvalget (Rogaland fylkeskommune) - kl. 11:30",
        "   • Formannskapet (Bjerkreim) - kl. 17:00",
        "",
        "📅 Torsdag 29. januar:",
        "   • Kommunestyret (Sirdal) - kl. 18:00",
        "   • Levekårsutvalget (Lund) - kl. 18:00"
    ])
    
    # 11. Postliste-bot
    add_content_slide(prs, "Postliste-bot", [
        "Skraping av postlister fra:",
        "   – Kommuner i dekningsområdet",
        "   – Rogaland fylkeskommune",
        "   – Ferde og andre aktører",
        "Nedlasting til Slack for videre bearbeiding",
        "Analyse med KI (Gemini, ChatGPT)"
    ])
    
    # 12. Brønnøysund-data
    add_content_slide(prs, "Data fra Brønnøysund", [
        "Uthenting av næringslivsdata",
        "Grunnlag for bransjeoversikter",
        "Restauranter, bedrifter, etc.",
        "Under konstruksjon – kjøres lokalt",
        "Potensial for automatiserte oppdateringer"
    ])
    
    # 13. Seksjonsslide - Teknisk
    add_section_slide(prs, "⚙️ Teknisk løsning")
    
    # 14. Integrasjoner
    add_content_slide(prs, "Integrasjoner i mm.schibsted.media", [
        "Embed-kode som i Mapcreator/Datawrapper",
        "KI genererer inline HTML",
        "Mobile-first design",
        "Tilpasset Aftenbladets blå designprofil",
        "Støtte for norske tegn (æ/ø/å)"
    ])
    
    # 15. LiteLLM
    add_content_slide(prs, "Tilgang til KI-modeller: LiteLLM", [
        "Open-source Python-bibliotek",
        "Unified interface for LLM-er",
        "Støtter mange modeller (GPT, Claude, etc.)",
        "Ca. 170 brukere i Schibsted Slack-kanal",
        "Be om tilgang via #liteLLM på Slack"
    ])
    
    # 16. Backup-løsninger
    add_content_slide(prs, "Når systemene svikter", [
        "Jojo nede med blank skjerm?",
        "Lag din egen transkribent!",
        "Uavhengighet fra enkeltverktøy",
        "Vibe koding gir fleksibilitet",
        "Fra problem til løsning på timer, ikke uker"
    ])
    
    # 17. Rutekutt-eksempel
    add_content_slide(prs, "Eksempel: Rutekutt", [
        "Fra regneark til interaktiv løsning",
        "Inline HTML-integrasjon",
        "Flertallspartienes forslag visualisert",
        "Lesbar på mobil og desktop",
        "Oppdateres enkelt ved behov"
    ], "Rutekutt v3-flertallspartiene")
    
    # 18. Oppsummering
    add_content_slide(prs, "Oppsummering", [
        "Vibe koding = KI + kreativitet + journalistikk",
        "Raskere fra idé til publisering",
        "Automatisering sparer tid",
        "Bedre datagrunnlag for saker",
        "Alle kan bidra – ikke bare utviklere"
    ])
    
    # 19. Avslutningsslide
    add_title_slide(prs, "Takk for oppmerksomheten!", "Spørsmål?\n\n#liteLLM på Slack for å komme i gang")
    
    # Lagre presentasjonen
    output_path = os.path.join(os.path.dirname(__file__), "vibe_koding_presentasjon.pptx")
    prs.save(output_path)
    print(f"✅ Presentasjon lagret: {output_path}")
    return output_path

if __name__ == "__main__":
    create_presentation()
